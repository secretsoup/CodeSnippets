import subprocess
import json
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches

# List of subscription IDs to review against
subscription_ids = ["sub-id-1", "sub-id-2", "sub-id-3"]  # Add your subscription IDs here

# Time period for the report (currently set to the previous month and the month before)
end_date = datetime.now().replace(day=1) - timedelta(days=1)  # Last day of the previous month
start_date_prev_month = end_date.replace(day=1)  # First day of the previous month
start_date_prev_prev_month = (start_date_prev_month - timedelta(days=1)).replace(day=1)  # First day of the month before

prs = Presentation()

def get_cost_for_subscription(sub_id, start_date, end_date): # TODO: Move this to Graph API call instead
    """Fetch cost data for a subscription using Azure CLI."""
    command = [
        "az", "consumption", "usage", "list",
        "--subscription", sub_id,
        "--start-date", start_date.strftime("%Y-%m-%d"),
        "--end-date", end_date.strftime("%Y-%m-%d"),
        "--output", "json"
    ]
    result = subprocess.run(command, capture_output=True, text=True)
    if result.returncode != 0:
        raise Exception(f"Failed to fetch cost data for subscription {sub_id}: {result.stderr}")
    return json.loads(result.stdout)

def calculate_total_cost(cost_data):
    """Calculate total cost from the cost data."""
    return sum(float(item["pretaxCost"]) for item in cost_data)

for sub_id in subscription_ids:
    try:
        # Fetch cost data for the previous month
        cost_data_prev_month = get_cost_for_subscription(sub_id, start_date_prev_month, end_date)
        total_cost_prev_month = calculate_total_cost(cost_data_prev_month)

        # Fetch cost data for the month before
        cost_data_prev_prev_month = get_cost_for_subscription(sub_id, start_date_prev_prev_month, start_date_prev_month - timedelta(days=1))
        total_cost_prev_prev_month = calculate_total_cost(cost_data_prev_prev_month)

        # Calculate cost difference
        cost_difference = total_cost_prev_month - total_cost_prev_prev_month
        cost_difference_percent = (cost_difference / total_cost_prev_prev_month) * 100 if total_cost_prev_prev_month != 0 else 0

        # Create a slide for the subscription
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"Subscription: {sub_id}"

        # Add cost data to the slide
        content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
        text_frame = content.text_frame
        text_frame.text = f"Cost for {start_date_prev_month.strftime('%B %Y')}: ${total_cost_prev_month:.2f}\n"
        text_frame.text += f"Cost for {start_date_prev_prev_month.strftime('%B %Y')}: ${total_cost_prev_prev_month:.2f}\n"
        text_frame.text += f"Cost Difference: ${cost_difference:.2f} ({cost_difference_percent:.2f}%)"

    except Exception as e:
        print(f"Error processing subscription {sub_id}: {e}")

# Save the PowerPoint presentation
prs.save("Cost_report.pptx")
print("PowerPoint report generated successfully!")
